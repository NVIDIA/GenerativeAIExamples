# SPDX-FileCopyrightText: Copyright (c) 2023-2024 NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
# http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import os
import subprocess

from langchain.docstore.document import Document
from pdfplumber import open as pdf_open
from pptx import Presentation

from RAG.examples.advanced_rag.multimodal_rag.vectorstore.custom_pdf_parser import is_graph, process_graph


def convert_ppt_to_pdf(ppt_path):
    """Convert a PowerPoint file to PDF using LibreOffice and save in '../../ppt_references/' folder."""
    base_name = os.path.basename(ppt_path)
    ppt_name_without_ext = os.path.splitext(base_name)[0].replace(" ", "_")

    # Use the existing directory '../../ppt_references/'
    new_dir_path = os.path.abspath("multimodal/ppt_references")

    # Set the new PDF path in the existing directory
    pdf_path = os.path.join(new_dir_path, f"{ppt_name_without_ext}.pdf")

    # LibreOffice command to convert PPT to PDF
    command = ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', new_dir_path, ppt_path]
    subprocess.run(command, check=True)

    return pdf_path


def convert_pdf_to_images(pdf_path):
    """Convert a PDF file to a series of images using PyMuPDF and save in '../../ppt_references/' folder."""
    doc = pdf_open(pdf_path)

    # Extract the base name of the PDF file and replace spaces with underscores
    base_name = os.path.basename(pdf_path)
    pdf_name_without_ext = os.path.splitext(base_name)[0].replace(" ", "_")

    # Use the existing directory '../../ppt_references/'
    new_dir_path = os.path.join(os.getcwd(), "multimodal/ppt_references")

    image_paths = []

    for page_num, page in enumerate(doc.pages):
        pix = page.to_image()

        # Save images in the existing directory
        output_image_path = os.path.join(new_dir_path, f"{pdf_name_without_ext}_{page_num:04d}.png")
        pix.save(output_image_path)
        image_paths.append((output_image_path, page_num))

    doc.close()
    return image_paths


def extract_text_and_notes_from_ppt(ppt_path):
    """Extract text and notes from a PowerPoint file."""
    prs = Presentation(ppt_path)
    text_and_notes = []
    for slide in prs.slides:
        slide_text = ' '.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        try:
            notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ""
        except:
            notes = ""
        text_and_notes.append((slide_text, notes))
    return text_and_notes


def process_ppt_file(ppt_path):
    """Process a PowerPoint file."""
    pdf_path = os.path.join(
        os.getcwd(),
        "multimodal/ppt_references",
        os.path.basename(ppt_path).replace('.pptx', '.pdf').replace('.ppt', '.pdf'),
    )
    convert_ppt_to_pdf(ppt_path)
    images_data = convert_pdf_to_images(pdf_path)
    slide_texts = extract_text_and_notes_from_ppt(ppt_path)
    processed_data = []

    for (image_path, page_num), (slide_text, notes) in zip(images_data, slide_texts):
        if notes:
            notes = "\n\nThe speaker notes for this slide are: " + notes

        # get image description with NeVA/DePlot
        image_description = " "
        if is_graph(image_path):
            image_description = process_graph(image_path)
        caption = slide_text + image_description + notes
        image_metadata = {
            "x1": 0,
            "y1": 0,
            "x2": 0,
            "x3": 0,
            "source": f"{os.path.basename(ppt_path)}",
            "image": image_path,
            "caption": caption,
            "type": "image",
            "page_num": page_num,
        }
        processed_data.append(
            Document(
                page_content="This is a slide with the text: " + slide_text + image_description,
                metadata=image_metadata,
            )
        )

    return processed_data
