# SPDX-FileCopyrightText: Copyright (c) 2023-2024 NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
# http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import os
import fitz
from pptx import Presentation
import subprocess
from llama_index.core import Document
from utils import (
    describe_image, is_graph, process_graph, extract_text_around_item, 
    process_text_blocks, save_uploaded_file
)

def get_pdf_documents(pdf_file):
    """Process a PDF file and extract text, tables, and images."""
    all_pdf_documents = []
    ongoing_tables = {}

    try:
        f = fitz.open(stream=pdf_file.read(), filetype="pdf")
    except Exception as e:
        print(f"Error opening or processing the PDF file: {e}")
        return []

    for i in range(len(f)):
        page = f[i]
        text_blocks = [block for block in page.get_text("blocks", sort=True) 
                       if block[-1] == 0 and not (block[1] < page.rect.height * 0.1 or block[3] > page.rect.height * 0.9)]
        grouped_text_blocks = process_text_blocks(text_blocks)
        
        table_docs, table_bboxes, ongoing_tables = parse_all_tables(pdf_file.name, page, i, text_blocks, ongoing_tables)
        all_pdf_documents.extend(table_docs)

        image_docs = parse_all_images(pdf_file.name, page, i, text_blocks)
        all_pdf_documents.extend(image_docs)

        for text_block_ctr, (heading_block, content) in enumerate(grouped_text_blocks, 1):
            heading_bbox = fitz.Rect(heading_block[:4])
            if not any(heading_bbox.intersects(table_bbox) for table_bbox in table_bboxes):
                bbox = {"x1": heading_block[0], "y1": heading_block[1], "x2": heading_block[2], "x3": heading_block[3]}
                text_doc = Document(
                    text=f"{heading_block[4]}\n{content}",
                    metadata={
                        **bbox,
                        "type": "text",
                        "page_num": i,
                        "source": f"{pdf_file.name[:-4]}-page{i}-block{text_block_ctr}"
                    },
                    id_=f"{pdf_file.name[:-4]}-page{i}-block{text_block_ctr}"
                )
                all_pdf_documents.append(text_doc)

    f.close()
    return all_pdf_documents

def parse_all_tables(filename, page, pagenum, text_blocks, ongoing_tables):
    """Extract tables from a PDF page."""
    table_docs = []
    table_bboxes = []
    try:
        tables = page.find_tables(horizontal_strategy="lines_strict", vertical_strategy="lines_strict")
        for tab in tables:
            if not tab.header.external:
                pandas_df = tab.to_pandas()
                tablerefdir = os.path.join(os.getcwd(), "vectorstore/table_references")
                os.makedirs(tablerefdir, exist_ok=True)
                df_xlsx_path = os.path.join(tablerefdir, f"table{len(table_docs)+1}-page{pagenum}.xlsx")
                pandas_df.to_excel(df_xlsx_path)
                bbox = fitz.Rect(tab.bbox)
                table_bboxes.append(bbox)

                before_text, after_text = extract_text_around_item(text_blocks, bbox, page.rect.height)

                table_img = page.get_pixmap(clip=bbox)
                table_img_path = os.path.join(tablerefdir, f"table{len(table_docs)+1}-page{pagenum}.jpg")
                table_img.save(table_img_path)
                description = process_graph(table_img.tobytes())

                caption = before_text.replace("\n", " ") + description + after_text.replace("\n", " ")
                if before_text == "" and after_text == "":
                    caption = " ".join(tab.header.names)
                table_metadata = {
                    "source": f"{filename[:-4]}-page{pagenum}-table{len(table_docs)+1}",
                    "dataframe": df_xlsx_path,
                    "image": table_img_path,
                    "caption": caption,
                    "type": "table",
                    "page_num": pagenum
                }
                all_cols = ", ".join(list(pandas_df.columns.values))
                doc = Document(text=f"This is a table with the caption: {caption}\nThe columns are {all_cols}", metadata=table_metadata)
                table_docs.append(doc)
    except Exception as e:
        print(f"Error during table extraction: {e}")
    return table_docs, table_bboxes, ongoing_tables

def parse_all_images(filename, page, pagenum, text_blocks):
    """Extract images from a PDF page."""
    image_docs = []
    image_info_list = page.get_image_info(xrefs=True)
    page_rect = page.rect

    for image_info in image_info_list:
        xref = image_info['xref']
        if xref == 0:
            continue

        img_bbox = fitz.Rect(image_info['bbox'])
        if img_bbox.width < page_rect.width / 20 or img_bbox.height < page_rect.height / 20:
            continue

        extracted_image = page.parent.extract_image(xref)
        image_data = extracted_image["image"]
        imgrefpath = os.path.join(os.getcwd(), "vectorstore/image_references")
        os.makedirs(imgrefpath, exist_ok=True)
        image_path = os.path.join(imgrefpath, f"image{xref}-page{pagenum}.png")
        with open(image_path, "wb") as img_file:
            img_file.write(image_data)

        before_text, after_text = extract_text_around_item(text_blocks, img_bbox, page.rect.height)
        if before_text == "" and after_text == "":
            continue

        image_description = " "
        if is_graph(image_data):
            image_description = process_graph(image_data)

        caption = before_text.replace("\n", " ") + image_description + after_text.replace("\n", " ")

        image_metadata = {
            "source": f"{filename[:-4]}-page{pagenum}-image{xref}",
            "image": image_path,
            "caption": caption,
            "type": "image",
            "page_num": pagenum
        }
        image_docs.append(Document(text="This is an image with the caption: " + caption, metadata=image_metadata))
    return image_docs

def process_ppt_file(ppt_path):
    """Process a PowerPoint file."""
    pdf_path = convert_ppt_to_pdf(ppt_path)
    images_data = convert_pdf_to_images(pdf_path)
    slide_texts = extract_text_and_notes_from_ppt(ppt_path)
    processed_data = []

    for (image_path, page_num), (slide_text, notes) in zip(images_data, slide_texts):
        if notes:
            notes = "\n\nThe speaker notes for this slide are: " + notes
        
        with open(image_path, 'rb') as image_file:
            image_content = image_file.read()
        
        image_description = " "
        if is_graph(image_content):
            image_description = process_graph(image_content)
        
        image_metadata = {
            "source": f"{os.path.basename(ppt_path)}",
            "image": image_path,
            "caption": slide_text + image_description + notes,
            "type": "image",
            "page_num": page_num
        }
        processed_data.append(Document(text="This is a slide with the text: " + slide_text + image_description, metadata=image_metadata))

    return processed_data

def convert_ppt_to_pdf(ppt_path):
    """Convert a PowerPoint file to PDF using LibreOffice."""
    base_name = os.path.basename(ppt_path)
    ppt_name_without_ext = os.path.splitext(base_name)[0].replace(' ', '_')
    new_dir_path = os.path.abspath("vectorstore/ppt_references")
    os.makedirs(new_dir_path, exist_ok=True)
    pdf_path = os.path.join(new_dir_path, f"{ppt_name_without_ext}.pdf")
    command = ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', new_dir_path, ppt_path]
    subprocess.run(command, check=True)
    return pdf_path

def convert_pdf_to_images(pdf_path):
    """Convert a PDF file to a series of images using PyMuPDF."""
    doc = fitz.open(pdf_path)
    base_name = os.path.basename(pdf_path)
    pdf_name_without_ext = os.path.splitext(base_name)[0].replace(' ', '_')
    new_dir_path = os.path.join(os.getcwd(), "vectorstore/ppt_references")
    os.makedirs(new_dir_path, exist_ok=True)
    image_paths = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap()
        output_image_path = os.path.join(new_dir_path, f"{pdf_name_without_ext}_{page_num:04d}.png")
        pix.save(output_image_path)
        image_paths.append((output_image_path, page_num))
    doc.close()
    return image_paths

def extract_text_and_notes_from_ppt(ppt_path):
    """Extract text and notes from a PowerPoint file."""
    prs = Presentation(ppt_path)
    text_and_notes = []
    for slide in prs.slides:
        slide_text = ' '.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        try:
            notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ''
        except:
            notes = ''
        text_and_notes.append((slide_text, notes))
    return text_and_notes

def load_multimodal_data(files):
    """Load and process multiple file types."""
    documents = []
    for file in files:
        file_extension = os.path.splitext(file.name.lower())[1]
        if file_extension in ('.png', '.jpg', '.jpeg'):
            image_content = file.read()
            image_text = describe_image(image_content)
            doc = Document(text=image_text, metadata={"source": file.name, "type": "image"})
            documents.append(doc)
        elif file_extension == '.pdf':
            try:
                pdf_documents = get_pdf_documents(file)
                documents.extend(pdf_documents)
            except Exception as e:
                print(f"Error processing PDF {file.name}: {e}")
        elif file_extension in ('.ppt', '.pptx'):
            try:
                ppt_documents = process_ppt_file(save_uploaded_file(file))
                documents.extend(ppt_documents)
            except Exception as e:
                print(f"Error processing PPT {file.name}: {e}")
        else:
            text = file.read().decode("utf-8")
            doc = Document(text=text, metadata={"source": file.name, "type": "text"})
            documents.append(doc)
    return documents

def load_data_from_directory(directory):
    """Load and process multiple file types from a directory."""
    documents = []
    for filename in os.listdir(directory):
        filepath = os.path.join(directory, filename)
        file_extension = os.path.splitext(filename.lower())[1]
        print(filename)
        if file_extension in ('.png', '.jpg', '.jpeg'):
            with open(filepath, "rb") as image_file:
                image_content = image_file.read()
            image_text = describe_image(image_content)
            doc = Document(text=image_text, metadata={"source": filename, "type": "image"})
            print(doc)
            documents.append(doc)
        elif file_extension == '.pdf':
            with open(filepath, "rb") as pdf_file:
                try:
                    pdf_documents = get_pdf_documents(pdf_file)
                    documents.extend(pdf_documents)
                except Exception as e:
                    print(f"Error processing PDF {filename}: {e}")
        elif file_extension in ('.ppt', '.pptx'):
            try:
                ppt_documents = process_ppt_file(filepath)
                documents.extend(ppt_documents)
                print(ppt_documents)
            except Exception as e:
                print(f"Error processing PPT {filename}: {e}")
        else:
            with open(filepath, "r", encoding="utf-8") as text_file:
                text = text_file.read()
            doc = Document(text=text, metadata={"source": filename, "type": "text"})
            documents.append(doc)
    return documents